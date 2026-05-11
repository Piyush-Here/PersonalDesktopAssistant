from __future__ import annotations

from pathlib import Path

from app.models.schemas import ExecutionResult

# ── Optional heavy imports guarded at module level ────────────────────────────
# Importing these unconditionally crashes the entire app on startup if any one
# library is missing. We import lazily inside each branch instead.

try:
    from docx import Document as _DocxDocument
    _DOCX_OK = True
except ImportError:
    _DOCX_OK = False

try:
    from openpyxl import load_workbook as _load_workbook
    _XLSX_OK = True
except ImportError:
    _XLSX_OK = False

try:
    from pptx import Presentation as _Presentation
    _PPTX_OK = True
except ImportError:
    _PPTX_OK = False

try:
    import fitz as _fitz          # PyMuPDF
    _PDF_ENGINE = "pymupdf"
except ImportError:
    try:
        from pypdf import PdfReader as _PdfReader
        _PDF_ENGINE = "pypdf"
    except ImportError:
        _PDF_ENGINE = None


class DocumentReaderTool:
    def preview(self, target: str) -> ExecutionResult:
        candidate = Path(target).expanduser()
        if candidate.exists() and candidate.is_file():
            return self._read_file(candidate)

        return ExecutionResult(
            success=True,
            message=f"No concrete file path resolved for '{target}'.",
            details=["Provide a full path or ask the assistant to find the file first."],
        )

    def _read_file(self, path: Path) -> ExecutionResult:
        suffix = path.suffix.lower()
        try:
            if suffix in {".txt", ".md", ".py", ".json"}:
                text = path.read_text(encoding="utf-8", errors="ignore")
                return self._summarize_text(path, text)

            if suffix == ".csv":
                text = path.read_text(encoding="utf-8", errors="ignore")
                return self._summarize_text(path, text)

            if suffix == ".pdf":
                return self._read_pdf(path)

            if suffix == ".docx":
                return self._read_docx(path)

            if suffix == ".pptx":
                return self._read_pptx(path)

            if suffix in {".xlsx", ".xlsm"}:
                return self._read_xlsx(path)

        except Exception as exc:
            return ExecutionResult(
                success=False,
                message=f"Failed to read {path.name}: {exc}",
                details=[str(exc)],
            )

        return ExecutionResult(
            success=True,
            message=f"Unsupported file type for {path.name}.",
            details=["Supported: txt, md, pdf, docx, pptx, xlsx, csv"],
        )

    # ── Per-format readers ────────────────────────────────────────────────────

    def _read_pdf(self, path: Path) -> ExecutionResult:
        if _PDF_ENGINE == "pymupdf":
            import fitz
            doc = fitz.open(str(path))
            text = "\n".join(doc[i].get_text() for i in range(min(5, len(doc))))
            doc.close()
            return self._summarize_text(path, text)
        if _PDF_ENGINE == "pypdf":
            from pypdf import PdfReader
            reader = PdfReader(str(path))
            text = "\n".join(page.extract_text() or "" for page in reader.pages[:5])
            return self._summarize_text(path, text)
        return ExecutionResult(
            success=False,
            message=f"Cannot read {path.name}: no PDF library installed.",
            details=["Run: pip install pymupdf"],
        )

    def _read_docx(self, path: Path) -> ExecutionResult:
        if not _DOCX_OK:
            return ExecutionResult(
                success=False,
                message=f"Cannot read {path.name}: python-docx is not installed.",
                details=["Run: pip install python-docx"],
            )
        from docx import Document
        doc = Document(str(path))
        text = "\n".join(p.text for p in doc.paragraphs[:50])
        return self._summarize_text(path, text)

    def _read_pptx(self, path: Path) -> ExecutionResult:
        if not _PPTX_OK:
            return ExecutionResult(
                success=False,
                message=f"Cannot read {path.name}: python-pptx is not installed.",
                details=["Run: pip install python-pptx"],
            )
        from pptx import Presentation
        prs = Presentation(str(path))
        slides = []
        for slide in prs.slides[:10]:
            fragments = [
                shape.text
                for shape in slide.shapes
                if hasattr(shape, "text") and shape.text
            ]
            slides.append(" ".join(fragments))
        return self._summarize_text(path, "\n".join(slides))

    def _read_xlsx(self, path: Path) -> ExecutionResult:
        if not _XLSX_OK:
            return ExecutionResult(
                success=False,
                message=f"Cannot read {path.name}: openpyxl is not installed.",
                details=["Run: pip install openpyxl"],
            )
        from openpyxl import load_workbook
        workbook = load_workbook(filename=str(path), read_only=True, data_only=True)
        details = []
        for sheet in workbook.sheetnames[:3]:
            ws = workbook[sheet]
            rows = list(ws.iter_rows(min_row=1, max_row=5, values_only=True))
            details.append(f"Sheet '{sheet}': {rows}")
        workbook.close()
        return ExecutionResult(
            success=True,
            message=f"Read workbook {path.name}.",
            details=details,
        )

    # ── Helpers ───────────────────────────────────────────────────────────────

    def _summarize_text(self, path: Path, text: str) -> ExecutionResult:
        normalized = " ".join(text.split())
        preview = normalized[:600] if normalized else "No readable text found."
        return ExecutionResult(
            success=True,
            message=f"Read {path.name}.",
            details=[preview],
        )
